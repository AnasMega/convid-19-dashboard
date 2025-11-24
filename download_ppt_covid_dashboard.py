from pptx import Presentation
from pptx.util import Inches

# Initialize Presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "COVID-19 Data Visualization Dashboard"
subtitle.text = "Data Visualization Project by Anas Hussain\nDate: 2025\nContact: anashussain0311@gmail.com"

# Slide 2: Introduction
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.shapes.placeholders[1]
title.text = "Introduction"
content.text = """
Overview of the project:
- Visualizing COVID-19 data from different countries
- Key components: Data cleaning, visualization, prediction models
- Tools: Python, Pandas, Seaborn, Matplotlib, Plotly, Streamlit
"""

# Slide 3: Dataset Overview
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.shapes.placeholders[1]
title.text = "Dataset Overview"
content.text = """
Data source: "Our World in Data" for COVID-19 statistics
Columns: Country, Date, DailyCasesPerMillion
Example of raw data will be shown.
"""

# Add more slides following the same pattern...

# Save the presentation
prs.save('covid_dashboard_presentation.pptx')
