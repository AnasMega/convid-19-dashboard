from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# Title Slide
slide_0 = prs.slides.add_slide(title_slide_layout)
title = slide_0.shapes.title
subtitle = slide_0.placeholders[1]
title.text = "COVID-19 Data Visualization Dashboard"
subtitle.text = "Developed by Anas Hussain\nDepartment of Computer Science\nData Science Institute, NED University"

# Slide 1: Overview
slide_1 = prs.slides.add_slide(bullet_slide_layout)
slide_1.shapes.title.text = "Project Overview"
content = slide_1.shapes.placeholders[1]
content.text = (
    "• A web-based interactive COVID-19 Dashboard using Streamlit\n"
    "• Shows global data on daily cases per million\n"
    "• Visualizations: Line Chart, Bar Chart, Heatmap, Choropleth Map, Animated Trend\n"
    "• Filter options: Country and Date Range\n"
    "• Tools: Python, Pandas, Seaborn, Plotly, Matplotlib, Streamlit"
)

# Slide 2: System Diagram
slide_2 = prs.slides.add_slide(prs.slide_layouts[5])
title = slide_2.shapes.title
title.text = "System Architecture Diagram"

img_path = "/mnt/data/covid_dashboard_diagram.png"
import matplotlib.pyplot as plt

# Create a simple system diagram
fig, ax = plt.subplots(figsize=(10, 6))
ax.axis('off')
components = [
    "CSV Dataset (Our World in Data)",
    "Data Processing with Pandas",
    "Visualizations (Seaborn, Matplotlib, Plotly)",
    "Streamlit Web App",
    "User Interface (Filters, Charts, Maps)"
]
positions = [(0.1, 0.8), (0.35, 0.6), (0.6, 0.6), (0.35, 0.4), (0.35, 0.2)]
for pos, comp in zip(positions, components):
    ax.add_patch(plt.Rectangle(pos, 0.25, 0.1, fill=True, color='lightblue', edgecolor='black'))
    ax.text(pos[0] + 0.125, pos[1] + 0.05, comp, ha='center', va='center', fontsize=10)
# Arrows
ax.annotate("", xy=(0.35, 0.75), xytext=(0.225, 0.85), arrowprops=dict(arrowstyle="->"))
ax.annotate("", xy=(0.6, 0.65), xytext=(0.475, 0.65), arrowprops=dict(arrowstyle="->"))
ax.annotate("", xy=(0.475, 0.5), xytext=(0.475, 0.6), arrowprops=dict(arrowstyle="->"))
ax.annotate("", xy=(0.475, 0.3), xytext=(0.475, 0.4), arrowprops=dict(arrowstyle="->"))
fig.savefig(img_path, bbox_inches='tight')
plt.close()

left = Inches(1.5)
top = Inches(1.8)
height = Inches(4.5)
slide_2.shapes.add_picture(img_path, left, top, height=height)

# Slide 3: Key Features
slide_3 = prs.slides.add_slide(bullet_slide_layout)
slide_3.shapes.title.text = "Key Features"
content = slide_3.shapes.placeholders[1]
content.text = (
    "• Multi-country and date range filtering\n"
    "• Animated line charts with Plotly\n"
    "• Weekly Heatmap using Seaborn\n"
    "• Downloadable filtered data CSV\n"
    "• Choropleth world map for selected dates\n"
    "• Real-time user interaction via Streamlit"
)

# Slide 4: Technologies Used
slide_4 = prs.slides.add_slide(bullet_slide_layout)
slide_4.shapes.title.text = "Technologies Used"
content = slide_4.shapes.placeholders[1]
content.text = (
    "• Python for backend logic\n"
    "• Pandas for data manipulation\n"
    "• Seaborn & Matplotlib for static visualizations\n"
    "• Plotly for interactive charts and map\n"
    "• Streamlit for dashboard UI\n"
    "• GitHub for version control"
)

# Slide 5: Conclusion & Contact
slide_5 = prs.slides.add_slide(bullet_slide_layout)
slide_5.shapes.title.text = "Conclusion & Contact"
content = slide_5.shapes.placeholders[1]
content.text = (
    "• A dynamic, extensible COVID-19 dashboard\n"
    "• Ideal for educational and analytical insights\n"
    "• Developed by: Anas Hussain\n"
    "• Email: anashussain0311@gmail.com\n"
    "• GitHub: github.com/anasMega\n"
    "• Institute: NED University of Engineering and Technology"
)

# Save presentation
pptx_path = "/Anas Hussain DV PPT FOr Covid 19 and Temperature Data Dashboard.pptx"
prs.save(pptx_path)

pptx_path
